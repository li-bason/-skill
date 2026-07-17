#!/usr/bin/env python3
"""Safely import study sources and maintain source_manifest.json."""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import shutil
import sys
import zipfile
from datetime import datetime, timezone
from pathlib import Path

TEXT = {".md", ".txt", ".csv", ".tsv", ".html", ".htm"}
IMAGES = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}


def digest(path: Path) -> str:
    h = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            h.update(block)
    return h.hexdigest()


def safe_zip(zip_path: Path, destination: Path) -> list[Path]:
    extracted: list[Path] = []
    with zipfile.ZipFile(zip_path) as archive:
        for member in archive.infolist():
            target = (destination / member.filename).resolve()
            try:
                target.relative_to(destination.resolve())
            except ValueError as exc:
                raise ValueError(f"ZIP 包含路径穿越：{member.filename}") from exc
            if member.is_dir():
                target.mkdir(parents=True, exist_ok=True)
                continue
            target.parent.mkdir(parents=True, exist_ok=True)
            with archive.open(member) as source, target.open("wb") as output:
                shutil.copyfileobj(source, output)
            extracted.append(target)
    return extracted


def extract_text(path: Path, output: Path) -> tuple[str, str | None]:
    ext = path.suffix.lower()
    try:
        if ext in TEXT:
            if ext in {".csv", ".tsv"}:
                delimiter = "\t" if ext == ".tsv" else ","
                with path.open(encoding="utf-8-sig", newline="") as handle:
                    rows = list(csv.reader(handle, delimiter=delimiter))
                text = "\n".join(" | ".join(row) for row in rows)
            else:
                text = path.read_text(encoding="utf-8", errors="replace")
        elif ext == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            text = "\n\n".join(
                f"## Page {i}\n{page.extract_text() or ''}" for i, page in enumerate(reader.pages, 1)
            )
        elif ext == ".docx":
            from docx import Document
            text = "\n".join(p.text for p in Document(str(path)).paragraphs)
        elif ext == ".pptx":
            from pptx import Presentation
            slides = []
            for i, slide in enumerate(Presentation(str(path)).slides, 1):
                chunks = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text]
                slides.append(f"## Slide {i}\n" + "\n".join(chunks))
            text = "\n\n".join(slides)
        elif ext == ".xlsx":
            from openpyxl import load_workbook
            workbook = load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sheet in workbook.worksheets:
                parts.append(f"## Sheet: {sheet.title}")
                parts.extend(" | ".join("" if v is None else str(v) for v in row) for row in sheet.iter_rows(values_only=True))
            text = "\n".join(parts)
        elif ext in IMAGES:
            return "needs_ocr", None
        else:
            return "needs_conversion", None
    except ImportError:
        return "needs_conversion", None
    except Exception as exc:
        return "failed", str(exc)
    # PDF text extractors occasionally return isolated UTF-16 surrogate code points.
    # Replace only unencodable fragments so the rest of the extracted material remains usable.
    text = text.encode("utf-8", errors="replace").decode("utf-8")
    output.write_text(text, encoding="utf-8")
    return "parsed", None


def import_one(task_dir: Path, source: Path, manifest: dict) -> list[dict]:
    if not source.is_file():
        raise ValueError(f"文件不存在：{source}")
    source_hash = digest(source)
    for item in manifest["sources"]:
        if item.get("sha256") == source_hash:
            return []
    source_id = f"SRC-{len(manifest['sources']) + 1:03d}"
    stored_name = f"{source_id}-{source.name}"
    stored = task_dir / "sources" / stored_name
    shutil.copy2(source, stored)
    if source.suffix.lower() == ".zip":
        extracted_dir = task_dir / "sources" / f"{source_id}-extracted"
        files = safe_zip(stored, extracted_dir)
        record = {
            "source_id": source_id, "original_name": source.name, "stored_path": stored.relative_to(task_dir).as_posix(),
            "sha256": source_hash, "status": "archive_extracted", "children": len(files)
        }
        manifest["sources"].append(record)
        records = [record]
        for child in files:
            records.extend(import_one(task_dir, child, manifest))
        return records
    normalized = task_dir / "sources" / f"{source_id}.md"
    status, error = extract_text(stored, normalized)
    record = {
        "source_id": source_id,
        "original_name": source.name,
        "stored_path": stored.relative_to(task_dir).as_posix(),
        "normalized_path": normalized.relative_to(task_dir).as_posix() if normalized.exists() else None,
        "sha256": source_hash,
        "status": status,
        "imported_at": datetime.now(timezone.utc).isoformat(),
        "error": error,
    }
    manifest["sources"].append(record)
    return [record]


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("task_dir")
    parser.add_argument("sources", nargs="+")
    args = parser.parse_args()
    task_dir = Path(args.task_dir).resolve()
    manifest_path = task_dir / "source_manifest.json"
    try:
        if not task_dir.is_dir():
            raise ValueError(f"任务目录不存在：{task_dir}")
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        imported = []
        for value in args.sources:
            imported.extend(import_one(task_dir, Path(value).resolve(), manifest))
        manifest_path.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    except (OSError, ValueError, json.JSONDecodeError, zipfile.BadZipFile) as exc:
        print(f"[import-sources] ERROR: {exc}", file=sys.stderr)
        return 1
    print(json.dumps({"imported": imported, "manifest": str(manifest_path)}, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
